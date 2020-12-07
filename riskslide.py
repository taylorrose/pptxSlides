from pptx import Presentation
import subprocess
from pandas import *
from riskgrid import *
import datetime
import os.path
from os import path

# TODO: Refactor as Classes
prs = Presentation('risk_master_template.pptx')


def xlscrape(path):
    """
    Convert Excel Sheet to python dictionary

    :param path: path where excel file resides
    :return: a dictionary of the first sheet in an excel file
    """
    xls = ExcelFile(path)
    return xls.parse(xls.sheet_names[0])


def print_placeholders(slide_target):
    """
    Used to see all placeholders on a slide with their index numbers, which are used for filling content in the PH

    :param slide_target: pass the specific slide object, e.g., prs._sides[0]
    :return: [the index number], [placeholder name] -- now that you know the placeholder index, you can use it to modify
            the placeholder, e.g., slide.placeholders[index].text = "hello, world"
    """
    for shape in slide_target.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('%d, %s' % (phf.idx, shape.name))


def Sort_Tuple(tup):
    """
    Function to sort the list of tuples by its second item

    :param tup: [('for', 24), ('is', 10), ('Geeks', 28),  ('Geeksforgeeks', 5), ('portal', 20), ('a', 15)]
    :return:[('Geeksforgeeks', 5), ('is', 10), ('a', 15), ('portal', 20), ('for', 24), ('Geeks', 28)]
    """
    lst = len(tup)
    for i in range(0, lst):

        for j in range(0, lst - i - 1):
            if (tup[j][1] > tup[j + 1][1]):
                temp = tup[j]
                tup[j] = tup[j + 1]
                tup[j + 1] = temp
    return tup.reverse()


rap = xlscrape("test.xlsx")

# TODO: Clean up data refinement mess

# ugly-ass list and dictionay comprehension to get
# [{'Incentivized Traffic':
#   {'Headline': ... ,
#   'Risk_Title': ... ,
#   'Ownership': ...},
#  {'Customer Life': ...}, ...
# ]
spreadsheet_data = [
    {rap["Risk_Title"][x]:
         {y: rap[y][x] for y in ["Headline",
                                 "Risk_Title",
                                 "Ownership",
                                 "Status",
                                 "Type",
                                 "Cause",
                                 "Effect",
                                 "Impact",
                                 "Contingency",
                                 "Mitigation",
                                 "Inherent_Severity",
                                 "Inherent_Likelihood",
                                 "Residual_Severity",
                                 "Residual_Likelihood",
                                 ]
          }
     }
    for x in range(len(rap))
]

obj = {}
for dic in spreadsheet_data:
    key = list(dic.keys())[0]
    obj[key] = dic[key]

for risk in list(obj.keys()):
    obj[risk]["Inherent_Rank"] = risk_score(obj[risk]["Inherent_Severity"], obj[risk]["Inherent_Likelihood"])
    obj[risk]["Residual_Rank"] = risk_score(obj[risk]["Residual_Severity"], obj[risk]["Residual_Likelihood"])
    obj[risk]["chart"] = str(obj[risk]["Inherent_Severity"]) + str(obj[risk]["Inherent_Likelihood"]) + str(
        obj[risk]["Residual_Severity"]) + str(obj[risk]["Residual_Likelihood"])

inherent_rank_tuple = [(obj[x]["Risk_Title"], int(str(obj[x]["Inherent_Rank"]) + str(obj[x]["Residual_Rank"]))) for x in
                       list(obj.keys())]
residual_rank_tuple = [(obj[x]["Risk_Title"], int(str(obj[x]["Residual_Rank"]) + str(obj[x]["Inherent_Rank"]))) for x in
                       list(obj.keys())]
Sort_Tuple(inherent_rank_tuple)
Sort_Tuple(residual_rank_tuple)
inherent_rank_tuple = [x[0] for x in inherent_rank_tuple]
inherent_rank_index = [list(rap["Risk_Title"]).index(x) for x in inherent_rank_tuple]
residual_rank_tuple = [x[0] for x in residual_rank_tuple]
residual_rank_index = [list(rap["Risk_Title"]).index(x) for x in residual_rank_tuple]


def print_grids_to_pngs():
    """
    Calls plot_risks function to create grids for all risks listed in spreadsheet. Rudimentary cache system checks if
    image already exists and bypasses process if png image is not needed, i.e., has already been created in folder

    :return: .png images stored in the img/risks/ directory. Images are stored by the inherity and residual severity
                plot points, e.g, 1211.png, 1312.png, etc.
    """
    for idx in inherent_rank_index:  # inherent_rank_index has index numbers each grid that is needed for presentation
        _ = [rap["Inherent_Severity"][idx],
             rap["Inherent_Likelihood"][idx],
             rap["Residual_Severity"][idx],
             rap["Residual_Likelihood"][idx]]
        chart = "".join([str(x) for x in _])  # combines all of _ list into string, e.g, if _ = [1,2,3,4] => "1234"

        if not path.exists("img/risks/" + chart + ".png"):
            plot_risk(_[0], _[1], _[2], _[3], chart)


print_grids_to_pngs()


def make_rap_slide(idx):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # rap = xlscrape("test.xlsx")

    slide_title = slide.placeholders[0]
    risk_grid = slide.placeholders[10]
    owner_tag = slide.placeholders[11]
    status_tag = slide.placeholders[12]
    type_tag = slide.placeholders[13]
    risk_title = slide.placeholders[14]
    cause_text = slide.placeholders[20]
    effect_text = slide.placeholders[22]
    impact_text = slide.placeholders[21]
    mitigation_text = slide.placeholders[23]
    contingency_text = slide.placeholders[24]

    slide_title.text = str(rap["Headline"][idx])
    risk_grid.insert_picture("img/risks/" + str(
        str(rap["Inherent_Severity"][idx]) + str(rap["Inherent_Likelihood"][idx]) + str(
            rap["Residual_Severity"][idx]) + str(rap["Residual_Likelihood"][idx])) + ".png")
    owner_tag.text = str(rap["Ownership"][idx])
    status_tag.text = str(rap["Status"][idx])
    type_tag.text = str(rap["Type"][idx])
    risk_title.text = str(rap["Risk_Title"][idx])

    cause_text.text = str(rap["Cause"][idx])
    effect_text.text = str(rap["Effect"][idx])
    impact_text.text = str(rap["Impact"][idx])
    mitigation_text.text = str(rap["Mitigation"][idx])
    contingency_text.text = str(rap["Contingency"][idx])


def make_title_slide(title, subhead=""):
    current_day = datetime.date.today()
    formatted_date = datetime.date.strftime(current_day, "%m.%d.%Y")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide_title = slide.placeholders[0]
    slide_subhead = slide.placeholders[10]
    slide_update = slide.placeholders[11]

    slide_title.text = title
    slide_subhead.text = subhead
    slide_update.text = "UPDATED: " + formatted_date


# TODO: Make Dynamic -- right now it will break if there aren't exactly 18 risks
# TODO: Make card part of image:
#   https://stackoverflow.com/questions/2563822/how-do-you-composite-an-image-onto-another-image-with-pil-in-python
def make_summary_slide(list):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.placeholders[0].text = "Risk Readiness"
    # if len(list) > 18:
    #     slide2 = prs.slides.add_slide(prs.slide_layouts[2])
    #     slide2.placeholders[0].text = "Risk Readiness cont."


    for x in range(len(list)):
        idx = inherent_rank_index[x]
        _ = [rap["Inherent_Severity"][idx],
             rap["Inherent_Likelihood"][idx],
             rap["Residual_Severity"][idx],
             rap["Residual_Likelihood"][idx]]
        chart = "".join([str(x) for x in _])  # combines all of _ list into string, e.g, if _ = [1,2,3,4] => "1234"

        if x < 18:
            slide.placeholders[x + 10].text = list[x]
            slide.placeholders[x + 28].insert_picture("img/risks/" + str(chart + ".png"))
            slide.placeholders[x + 46].text = "p." + str(x + 3)
            print("first sheet:", x)

        # if x >= 18:
        #     slide2.placeholders[x + 10].text = list[x]
        #     slide2.placeholders[x + 28].insert_picture("img/risks/" + str(chart + ".png"))
        #     slide2.placeholders[x + 46].text = "p." + str(x + 3)
        #     print("second sheet:", x)



make_title_slide("Risk Management Plan",
                 "An Example of a Dynamically Generated PowerPoint Slide Deck"
                 )

make_summary_slide(inherent_rank_tuple)

for x in inherent_rank_index:
    make_rap_slide(x)

prs.slides.add_slide(prs.slide_layouts[3])
prs.slides.add_slide(prs.slide_layouts[4])
prs.slides.add_slide(prs.slide_layouts[5])

FileName = "pptfile.pptx"
prs.save(FileName)

subprocess.call(['open', FileName])

## Convert date to MM/DD/YYYY format in Python
